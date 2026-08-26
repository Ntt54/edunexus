"""Text extraction + chunking for the local AI tutor (research D3).

UI-framework-free by contract (no textual/fastapi imports).
"""

from __future__ import annotations

import hashlib
import re
import zipfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable


SUPPORTED_FORMATS = {".txt", ".md", ".pdf", ".epub", ".docx", ".pptx"}


class _TagStripper(HTMLParser):
    """Collect visible text, dropping all markup (research D3 EPUB path)."""

    _BLOCK = {
        "p", "div", "br", "li", "tr", "td", "th", "section", "article",
        "blockquote", "h1", "h2", "h3", "h4", "h5", "h6",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        self._parts.append(data)

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in self._BLOCK:
            self._parts.append(" ")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._BLOCK:
            self._parts.append(" ")

    def get_text(self) -> str:
        return "".join(self._parts)

    def reset_text(self) -> None:
        self._parts = []


def extract_text(path, fmt: str | None = None) -> str:
    """Extract plain text from a book file.

    Args:
        path: Path to the document.
        fmt: Optional explicit format (``"txt"``, ``"md"``, ``"pdf"``,
            ``"epub"``); auto-detected from the file extension when omitted.

    Returns:
        The extracted plain-text content.

    Raises:
        FileNotFoundError: when the path does not exist.
        ValueError: for unsupported formats.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"No such file: {p}")
    suffix = p.suffix.lower()
    if fmt is None:
        if suffix not in SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {suffix or '(none)'}")
        fmt = suffix.lstrip(".")
    else:
        fmt = fmt.lower().lstrip(".")
        if fmt not in {f.lstrip(".") for f in SUPPORTED_FORMATS}:
            raise ValueError(f"Unsupported format: {fmt}")

    if fmt in ("txt", "md"):
        return p.read_text(encoding="utf-8", errors="replace")
    if fmt == "pdf":
        return _extract_pdf(p)
    if fmt == "epub":
        return _extract_epub(p)
    if fmt == "docx":
        return _extract_docx(p)
    if fmt == "pptx":
        return _extract_pptx(p)
    raise ValueError(f"Unsupported format: {fmt}")


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    """Extract text from a .docx file via python-docx."""
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from a .pptx file via python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n".join(parts)


_OPF_NS = "{http://www.idpf.org/2007/opf}"


def _find_opf(zf: zipfile.ZipFile, names: list[str]) -> str | None:
    if "META-INF/container.xml" in names:
        try:
            root = ET.fromstring(zf.read("META-INF/container.xml"))
        except ET.ParseError:
            root = None
        if root is not None:
            for el in root.iter():
                if el.tag.endswith("rootfile"):
                    fp = el.get("full-path")
                    if fp:
                        return fp
    for n in names:
        if n.endswith(".opf"):
            return n
    return None


def _spine_hrefs(zf: zipfile.ZipFile, opf_path: str) -> list[str]:
    try:
        root = ET.fromstring(zf.read(opf_path))
    except ET.ParseError:
        return []
    manifest: dict[str, str] = {}
    for item in root.iter(_OPF_NS + "item"):
        hid = item.get("id")
        href = item.get("href")
        if hid and href:
            manifest[hid] = href
    base = opf_path.rsplit("/", 1)[0] if "/" in opf_path else ""
    hrefs: list[str] = []
    for ref in root.iter(_OPF_NS + "itemref"):
        hid = ref.get("idref")
        if hid and hid in manifest:
            href = manifest[hid]
            hrefs.append(f"{base}/{href}" if base and not href.startswith("/") else href)
    return hrefs


def _extract_epub(path: Path) -> str:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        opf_path = _find_opf(zf, names)
        hrefs = _spine_hrefs(zf, opf_path) if opf_path else []
        if not hrefs:
            hrefs = [n for n in names if n.endswith((".xhtml", ".html", ".htm"))]
        stripper = _TagStripper()
        texts: list[str] = []
        for href in hrefs:
            try:
                data = zf.read(href).decode("utf-8", "replace")
            except KeyError:
                continue
            stripper.feed(data)
            texts.append(stripper.get_text())
            stripper.reset_text()
    return "\n".join(t for t in texts if t).strip()


def fingerprint(text: str) -> str:
    """sha256 of whitespace-normalized text (research D5)."""
    normalized = " ".join(text.split())
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


def chunk_text(
    text: str,
    max_chars: int = 1200,
    overlap: int = 200,
) -> list[str]:
    """Split ``text`` into word-boundary chunks (research D3).

    Each chunk holds at most ``max_chars`` characters (a single over-long word
    is kept whole). Consecutive chunks overlap by roughly ``overlap``
    characters, snapped to word boundaries so no word is ever split across a
    chunk boundary.
    """
    if max_chars <= 0:
        raise ValueError("max_chars must be positive")
    if overlap < 0:
        raise ValueError("overlap must be non-negative")
    if overlap >= max_chars:
        overlap = max_chars - 1

    words = text.split()
    if not words:
        return []

    chunks: list[str] = []
    start = 0
    n = len(words)
    while start < n:
        end = start
        total = 0
        while end < n:
            w = words[end]
            add = len(w) + (1 if end > start else 0)
            if total + add > max_chars and end > start:
                break
            total += add
            end += 1
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end >= n:
            break
        # step back for overlap, snapped to word boundaries
        step = 0
        covered = 0
        for w in reversed(words[start:end]):
            if covered >= overlap:
                break
            covered += len(w) + 1
            step += 1
        next_start = end - step
        if next_start <= start:
            next_start = start + 1  # guarantee forward progress
        start = next_start
    return chunks


# ----------------------------------------------------------------------
# Structured chunking (US4 — T027 / T028 / T029)
# ----------------------------------------------------------------------

_HEADING_RE = re.compile(r"^(#{1,4})\s+(.+)$", re.MULTILINE)
_PAGE_RE = re.compile(r"^(?:---\s*Page\s+(\d+)\s*---|Page\s+(\d+))$", re.MULTILINE)
_SENTENCE_SPLIT = re.compile(r"(?<=[.!?])\s+")


def chunk_text_structured(
    text: str,
    max_chars: int = 1200,
    overlap: int = 200,
    metadata: dict | None = None,
) -> list[dict]:
    """Split *text* into structured chunk dicts respecting paragraph boundaries.

    Each returned dict has the shape::

        {"text": str, "section": str | None, "page": int | None}

    Algorithm (T027):
      1. Pre-scan for headings (``^#{1,4} …``) and page markers
         (``--- Page X ---`` or ``Page X``) to build metadata.
      2. Split on ``\\n\\n`` (paragraph boundaries).  Heading paragraphs
         are merged with their following content paragraph.
      3. Individual paragraphs exceeding *max_chars* are sub-split at
         sentence boundaries, producing multiple chunks.
      4. Consecutive chunks share the last *overlap* characters.

    Args:
        text: Source text (typically Markdown).
        max_chars: Maximum character count per chunk (soft limit; a single
            over-long sentence is kept whole).
        overlap: Number of characters of overlap carried from one chunk
            to the next.
        metadata: Optional extra metadata merged into every chunk dict.

    Returns:
        A list of chunk dicts.  Empty input → empty list.
    """
    if max_chars <= 0:
        raise ValueError("max_chars doit être positif")
    if overlap < 0:
        raise ValueError("overlap doit être ≥ 0")
    if overlap >= max_chars:
        overlap = max_chars - 1

    if not text or not text.strip():
        return []

    base_meta: dict = dict(metadata) if metadata else {}

    # --- T028: heading scan ------------------------------------------------
    headings: list[tuple[int, str]] = []          # (char_offset, heading_text)
    for m in _HEADING_RE.finditer(text):
        headings.append((m.start(), m.group(2).strip()))

    # --- T029: page-number scan --------------------------------------------
    pages: list[tuple[int, int]] = []             # (char_offset, page_no)
    for m in _PAGE_RE.finditer(text):
        page_no = int(m.group(1) or m.group(2))
        pages.append((m.start(), page_no))

    # --- helpers to resolve current heading / page at a given offset ---------
    def _current_heading(offset: int) -> str | None:
        current: str | None = None
        for h_off, h_text in headings:
            if h_off <= offset:
                current = h_text
            else:
                break
        return current

    def _current_page(offset: int) -> int | None:
        current: int | None = None
        for p_off, p_no in pages:
            if p_off <= offset:
                current = p_no
            else:
                break
        return current

    def _is_heading(para: str) -> bool:
        return bool(_HEADING_RE.match(para.strip()))

    # --- split into paragraphs on \n\n ------------------------------------
    raw_paragraphs: list[str] = []
    for para in text.split("\n\n"):
        stripped = para.strip()
        if stripped:
            raw_paragraphs.append(stripped)

    if not raw_paragraphs:
        return []

    # --- merge heading paragraphs with their following content paragraph -----
    merged_paras: list[str] = []
    i = 0
    while i < len(raw_paragraphs):
        if _is_heading(raw_paragraphs[i]) and i + 1 < len(raw_paragraphs):
            # Merge heading + next content paragraph
            merged_paras.append(raw_paragraphs[i] + "\n\n" + raw_paragraphs[i + 1])
            i += 2
        else:
            merged_paras.append(raw_paragraphs[i])
            i += 1

    # --- resolve char offsets for heading/page lookup -----------------------
    para_offsets: list[int] = []
    off = 0
    for para in merged_paras:
        para_offsets.append(off)
        # Account for the original text offset of this paragraph
        # Find this paragraph's text in the original to get the right offset
        pos = text.find(para.split("\n\n")[0][:50], off)
        if pos >= 0:
            off = pos + len(para.split("\n\n")[0])
        else:
            off += len(para)

    # Simpler approach: compute offsets from the original paragraph list
    # to get correct heading/page resolution
    orig_offsets: list[int] = []
    off = 0
    for para in raw_paragraphs:
        orig_offsets.append(off)
        off += len(para) + 2  # +2 for \n\n

    # For merged paras (heading+content), use the heading's offset
    merged_offsets: list[int] = []
    idx = 0
    while idx < len(raw_paragraphs):
        if _is_heading(raw_paragraphs[idx]) and idx + 1 < len(raw_paragraphs):
            merged_offsets.append(orig_offsets[idx])
            idx += 2
        else:
            merged_offsets.append(orig_offsets[idx])
            idx += 1

    para_meta: list[tuple[str | None, int | None]] = [
        (_current_heading(o), _current_page(o)) for o in merged_offsets
    ]

    # --- sub-split long merged paragraphs at sentence boundaries ------------
    # Each resulting piece becomes its own chunk.
    # Track which paragraph each piece belongs to for offset resolution.
    final_texts: list[str] = []
    final_meta: list[tuple[str | None, int | None]] = []

    # Effective max for sub-split pieces leaves room for overlap carry
    effective_max = max_chars - overlap if overlap > 0 and max_chars > overlap else max_chars

    for pidx, para in enumerate(merged_paras):
        if len(para) <= max_chars:
            final_texts.append(para)
            final_meta.append(para_meta[pidx])
            continue
        # Sub-split at sentence boundaries
        sentences = _SENTENCE_SPLIT.split(para)
        buf = ""
        for sent in sentences:
            candidate = (buf + " " + sent).strip() if buf else sent
            if len(candidate) <= effective_max:
                buf = candidate
            elif not buf:
                # single sentence longer than effective_max → keep it whole
                final_texts.append(sent)
                final_meta.append(para_meta[pidx])
            else:
                final_texts.append(buf)
                final_meta.append(para_meta[pidx])
                buf = sent
        if buf:
            final_texts.append(buf)
            final_meta.append(para_meta[pidx])

    if not final_texts:
        return []

    # --- build chunk dicts with overlap ------------------------------------
    chunks: list[dict] = []
    carry: str = ""

    for idx, ft in enumerate(final_texts):
        if carry:
            effective = carry + "\n\n" + ft
        else:
            effective = ft

        h, p = final_meta[idx]
        chunk_entry: dict = {"text": effective}
        chunk_entry["section"] = base_meta.get("section", h)
        chunk_entry["page"] = base_meta.get("page", p)
        chunks.append(chunk_entry)

        # Overlap: tail of the final text carried to the next chunk
        if overlap > 0 and len(ft) > overlap:
            carry = ft[-overlap:]
        else:
            carry = ft

    return chunks
